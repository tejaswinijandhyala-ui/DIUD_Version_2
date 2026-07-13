"""
tools/export_builders.py — turns a stored query result or conversation
into a downloadable file. Ported from main.py's _build_csv, _build_xlsx,
_build_pdf, _build_pptx, and _generate_export_content.

Scope note: this port covers CSV, styled XLSX, and PDF/PPTX in SUMMARY
mode (an AI-generated executive summary). The original also had a
VERBATIM mode that re-embeds every chat turn's inline HTML charts as
screenshots via Playwright — that's a heavier dependency (a full
headless browser) and a more involved port, left out of this pass
deliberately rather than included half-working. Summary mode covers the
common case; verbatim mode is a clean follow-up if you want it.
"""

import csv
import io
import re
from datetime import date
from typing import Dict, List, Optional

from openai import OpenAI
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    BaseDocTemplate, Frame, PageTemplate, PageBreak,
    Paragraph, Spacer, Table, TableStyle,
)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

client = OpenAI()


def _safe_filename(title: str) -> str:
    return re.sub(r'[^\w\-]', '_', title)[:60]


def _strip_md(t: str) -> str:
    t = re.sub(r'\*{1,3}(.*?)\*{1,3}', r'\1', t)
    t = re.sub(r'#{1,6}\s*', '', t)
    t = re.sub(r'^[\-\*\u2022]\s*', '', t, flags=re.M)
    t = re.sub(r'`(.*?)`', r'\1', t)
    t = t.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    return t.strip()


# =============================================================================
# CSV
# =============================================================================

def build_csv(columns: List[str], rows: List[dict], sql: str = "", filters_applied: str = "") -> bytes:
    buf = io.StringIO()
    buf.write(f"# Generated: {date.today().isoformat()}\n")
    buf.write(f"# Total Records: {len(rows)}\n")
    buf.write(f"# Filters: {filters_applied}\n")
    buf.write("#\n")
    writer = csv.DictWriter(buf, fieldnames=columns, extrasaction="ignore", lineterminator="\n")
    writer.writeheader()
    for row in rows:
        writer.writerow(row)
    return buf.getvalue().encode("utf-8")


# =============================================================================
# XLSX -- styled Excel export
# =============================================================================

_CURRENCY_COLUMN_HINTS = ("amount", "value", "revenue", "target", "pipeline", "arr", "acv", "price", "cost", "spend", "budget", "quota")


def _looks_like_currency_column(col_name: str) -> bool:
    lower = col_name.lower()
    return any(hint in lower for hint in _CURRENCY_COLUMN_HINTS)


def build_xlsx(columns: List[str], rows: List[dict], title: str, filters_applied: str = "") -> bytes:
    """
    Bold white-on-navy header, currency-formatted amount/value/target
    columns, auto-sized columns, frozen header row. Always writes every
    row given to it -- if that count ever needs to differ from a claimed
    total, that mismatch should surface loudly at the call site, not be
    silently absorbed here.
    """
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"[:31]

    header_font = Font(name="Arial", size=11, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1E3A5F", end_color="1E3A5F", fill_type="solid")
    body_font = Font(name="Arial", size=10.5)
    header_align = Alignment(horizontal="center", vertical="center")

    for col_idx, col_name in enumerate(columns, start=1):
        cell = ws.cell(row=1, column=col_idx, value=col_name)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_align
    ws.freeze_panes = "A2"

    currency_cols = {c for c in columns if _looks_like_currency_column(c)}

    for row_idx, row in enumerate(rows, start=2):
        for col_idx, col_name in enumerate(columns, start=1):
            raw_val = row.get(col_name, "")
            cell = ws.cell(row=row_idx, column=col_idx)
            cell.font = body_font
            if col_name in currency_cols:
                try:
                    cell.value = float(raw_val)
                    cell.number_format = '$#,##0'
                    continue
                except (TypeError, ValueError):
                    pass
            cell.value = raw_val

    for col_idx, col_name in enumerate(columns, start=1):
        sample_lens = [len(str(r.get(col_name, ""))) for r in rows[:200]]
        width = max([len(col_name)] + sample_lens) + 3
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max(width, 10), 60)

    meta_row = len(rows) + 3
    for i, line in enumerate([
        f"Generated: {date.today().isoformat()}",
        f"Filters: {filters_applied or 'none'}",
        f"All {len(rows)} matching records included",
    ]):
        c = ws.cell(row=meta_row + i, column=1, value=line)
        c.font = Font(name="Arial", size=9, italic=True, color="666666")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# =============================================================================
# Content generation -- summary mode
# =============================================================================

def generate_export_summary(
    conversation: List[dict],
    title: str,
    export_type: str,
    columns: Optional[List[str]] = None,
    rows: Optional[List[dict]] = None,
) -> str:
    """
    Asks Claude for a concise executive summary of the conversation, in
    the shape appropriate for the target format. This is a separate,
    narrowly-scoped call from every agent in the live pipeline -- export
    generation isn't part of answering a question, it's repackaging
    answers that were already validated.
    """
    conv_text = "\n\n".join(
        f"{'USER' if m.get('role') == 'user' else 'ASSISTANT'}: {m.get('content', '')}"
        for m in conversation
    )
    format_hint = (
        "Format as a PowerPoint: use SLIDE: <title> for each slide, then bullet points."
        if export_type == "pptx"
        else "Format as a professional PDF report: ## section headers, narrative prose, tables."
    )
    dataset_hint = ""
    if rows:
        dataset_hint = f"\n\nDATASET CONTEXT: {len(rows)} records with columns: {', '.join((columns or [])[:12])}."

    prompt = f"""You are preparing a professional {export_type.upper()} summary report.

CONVERSATION:
{conv_text}
{dataset_hint}

TASK: Create a concise executive summary titled "{title}"

{format_hint}

REQUIREMENTS:
- Executive summary at the start with key numbers only
- Logical sections: summary, key metrics, insights, recommendations
- Bold key numbers; clean professional tone
- Today: {date.today().strftime('%B %d, %Y')}
- Generate the COMPLETE document -- do not truncate

Generate the summary report now:"""

    response = client.messages.create(
        model="gpt-5",
        system="You are a professional document formatter. Follow the instructions exactly. Never add content not requested.",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=6144,
    )
    return "\n".join(b.text for b in response.content if hasattr(b, "text") and b.text).strip()


def _rows_to_markdown_table(columns: List[str], rows: List[dict]) -> str:
    if not rows:
        return "_No data._"
    header = "| " + " | ".join(columns) + " |"
    sep = "| " + " | ".join("---" for _ in columns) + " |"
    lines = [header, sep]
    for row in rows:
        cells = [str(row.get(c, "")).replace("|", "\\|") for c in columns]
        lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


# =============================================================================
# PDF
# =============================================================================

_C_NAVY = colors.HexColor("#0D1B3E")
_C_BLUE = colors.HexColor("#1565C0")
_C_WHITE = colors.white
_C_BG = colors.HexColor("#F7F9FC")
_C_TXT = colors.HexColor("#1E293B")
_C_DIM = colors.HexColor("#94A3B8")

PW, PH = A4
_ML = _MR = 0.6 * inch
_MT = 0.45 * inch
_MB = 0.40 * inch
_HDR_H = 44
_FTR_H = 20
_CW = PW - _ML - _MR


def _pdf_styles():
    return {
        "Cover_Title": ParagraphStyle("Cover_Title", fontSize=26, leading=32, textColor=_C_WHITE, fontName="Helvetica-Bold", spaceAfter=8),
        "Cover_Sub": ParagraphStyle("Cover_Sub", fontSize=13, leading=18, textColor=colors.HexColor("#B0BEC5"), fontName="Helvetica"),
        "Section_H": ParagraphStyle("Section_H", fontSize=11, leading=15, textColor=_C_WHITE, fontName="Helvetica-Bold"),
        "Body": ParagraphStyle("Body", fontSize=9, leading=14, textColor=_C_TXT, fontName="Helvetica", spaceAfter=4),
        "Bullet": ParagraphStyle("Bullet", fontSize=9, leading=14, textColor=_C_TXT, fontName="Helvetica", leftIndent=12, firstLineIndent=-8, spaceAfter=3),
        "H2": ParagraphStyle("H2", fontSize=11, leading=15, textColor=_C_NAVY, fontName="Helvetica-Bold", spaceBefore=10, spaceAfter=4),
        "H3": ParagraphStyle("H3", fontSize=9, leading=13, textColor=_C_BLUE, fontName="Helvetica-Bold", spaceBefore=6, spaceAfter=2),
        "TH": ParagraphStyle("TH", fontSize=7, leading=9, textColor=_C_WHITE, fontName="Helvetica-Bold"),
        "TD": ParagraphStyle("TD", fontSize=7, leading=9, textColor=_C_TXT, fontName="Helvetica"),
    }


def _parse_sections(text: str):
    parts = re.split(r'^##\s+', text, flags=re.MULTILINE)
    return [
        (lines[0].strip(), lines[1].strip() if len(lines) > 1 else "")
        for part in parts if part.strip()
        for lines in [part.strip().split("\n", 1)]
    ]


def _md_table_to_rl(table_lines: list, styles: dict):
    data = []
    for idx, line in enumerate(table_lines):
        if "---" in line:
            continue
        cells = [c.strip().replace("\\|", "|") for c in line.strip("|").split("|")]
        row = [Paragraph(_strip_md(c), styles["TH"] if idx == 0 else styles["TD"]) for c in cells]
        data.append(row)
    if not data:
        return Spacer(1, 1)
    num_cols = max(len(r) for r in data)
    col_w = _CW / max(num_cols, 1)
    tbl = Table(data, colWidths=[col_w] * num_cols, repeatRows=1)
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), _C_NAVY),
        ("TEXTCOLOR", (0, 0), (-1, 0), _C_WHITE),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 7),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
        ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#E2E8F0")),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    return tbl


def build_pdf(title: str, report_text: str) -> bytes:
    buf = io.BytesIO()
    styles = _pdf_styles()

    def _on_page(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(_C_NAVY)
        canvas.rect(0, PH - _HDR_H - _MT, PW, _HDR_H + _MT, fill=1, stroke=0)
        canvas.setFillColor(_C_WHITE)
        canvas.setFont("Helvetica-Bold", 10)
        canvas.drawString(_ML, PH - _MT - 28, title)
        canvas.setFillColor(_C_BG)
        canvas.rect(0, 0, PW, _FTR_H + _MB, fill=1, stroke=0)
        canvas.setFillColor(_C_DIM)
        canvas.setFont("Helvetica", 7)
        canvas.drawCentredString(PW / 2, _MB + 5, f"AI-Generated | CONFIDENTIAL | {date.today().strftime('%B %Y')}")
        canvas.drawRightString(PW - _MR, _MB + 5, f"Page {canvas.getPageNumber()}")
        canvas.restoreState()

    frame = Frame(_ML, _MB + _FTR_H, _CW, PH - _HDR_H - _MT - _MB - _FTR_H, id="main")
    template = PageTemplate(id="main", frames=[frame], onPage=_on_page)
    doc = BaseDocTemplate(buf, pagesize=A4, leftMargin=_ML, rightMargin=_MR, topMargin=_MT + _HDR_H, bottomMargin=_MB + _FTR_H)
    doc.addPageTemplates([template])

    story = [
        Spacer(1, 1.0 * inch),
        Paragraph(title, styles["Cover_Title"]),
        Paragraph(f"Generated {date.today().strftime('%B %d, %Y')}", styles["Cover_Sub"]),
        PageBreak(),
    ]

    sections = _parse_sections(report_text)
    if not sections:
        for line in report_text.split("\n"):
            line = line.strip()
            if line:
                story.append(Paragraph(_strip_md(line), styles["Body"]))
    else:
        for sec_title, sec_body in sections:
            story.append(Table(
                [[Paragraph(sec_title.upper(), styles["Section_H"])]],
                colWidths=[_CW],
                style=TableStyle([
                    ("BACKGROUND", (0, 0), (-1, -1), _C_BLUE),
                    ("TOPPADDING", (0, 0), (-1, -1), 8),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                    ("LEFTPADDING", (0, 0), (-1, -1), 10),
                ])
            ))
            story.append(Spacer(1, 6))
            lines = sec_body.split("\n")
            i = 0
            while i < len(lines):
                line = lines[i].strip()
                if line.startswith("|") and i + 1 < len(lines) and "---" in lines[i + 1]:
                    table_lines = []
                    while i < len(lines) and lines[i].strip().startswith("|"):
                        table_lines.append(lines[i].strip())
                        i += 1
                    story.append(_md_table_to_rl(table_lines, styles))
                    story.append(Spacer(1, 6))
                    continue
                if not line:
                    story.append(Spacer(1, 3))
                elif line.startswith("### "):
                    story.append(Paragraph(_strip_md(line), styles["H3"]))
                elif line.startswith("## "):
                    story.append(Paragraph(_strip_md(line), styles["H2"]))
                elif line.startswith(("- ", "* ", "\u2022 ")):
                    story.append(Paragraph("\u2022 " + _strip_md(line[2:]), styles["Bullet"]))
                else:
                    story.append(Paragraph(_strip_md(line), styles["Body"]))
                i += 1
            story.extend([Spacer(1, 12), PageBreak()])

    doc.build(story)
    return buf.getvalue()


# =============================================================================
# PPTX
# =============================================================================

_C_NAVY_P = RGBColor(0x0D, 0x1B, 0x3E)
_C_DNAV_P = RGBColor(0x0A, 0x11, 0x28)
_C_BLUE_P = RGBColor(0x1E, 0x88, 0xE5)
_C_WHITE_P = RGBColor(0xFF, 0xFF, 0xFF)
_C_LTBG_P = RGBColor(0xF5, 0xF7, 0xFA)
_C_TXT_P = RGBColor(0x1A, 0x1A, 0x2E)
_C_DIM_P = RGBColor(0x88, 0x99, 0xAA)


def _pptx_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _pptx_rect(slide, l, t, w, h, color):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _pptx_txt(slide, text, l, t, w, h, bold=False, size=18, color=None, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or _C_TXT_P
    return txb


def _parse_slides(text: str):
    slides, cur_title, cur_bullets = [], None, []
    for line in text.split("\n"):
        line = line.rstrip()
        if line.startswith("SLIDE:"):
            if cur_title is not None:
                slides.append((cur_title, cur_bullets))
            cur_title, cur_bullets = line[6:].strip(), []
        elif line.startswith("- ") and cur_title:
            cur_bullets.append(line[2:].strip())
    if cur_title is not None:
        slides.append((cur_title, cur_bullets))
    return slides


def build_pptx(title: str, slide_text: str) -> bytes:
    slides_data = _parse_slides(slide_text) or [(title, [slide_text[:400]])]
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    footer_text = f"AI-Generated | CONFIDENTIAL | {date.today().strftime('%B %Y')}"
    blank = prs.slide_layouts[6]

    def _footer(s):
        _pptx_rect(s, 0, 7.1, 13.33, 0.4, _C_DNAV_P)
        _pptx_txt(s, footer_text, 0.3, 7.12, 12, 0.35, size=7, color=_C_DIM_P, align=PP_ALIGN.CENTER)

    cover = prs.slides.add_slide(blank)
    _pptx_bg(cover, _C_NAVY_P)
    _pptx_rect(cover, 0, 3.2, 13.33, 0.06, _C_BLUE_P)
    _pptx_txt(cover, title, 0.8, 1.6, 11.5, 1.4, bold=True, size=34, color=_C_WHITE_P)
    _pptx_txt(cover, f"Generated: {date.today().strftime('%B %d, %Y')}", 0.8, 3.6, 6, 0.45, size=12, color=RGBColor(0x78, 0x90, 0x9C))

    for i, (s_title, bullets) in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        _pptx_bg(slide, _C_LTBG_P)
        _pptx_rect(slide, 0, 0, 13.33, 0.9, _C_BLUE_P)
        _pptx_txt(slide, s_title.upper(), 0.35, 0.1, 12.5, 0.7, bold=True, size=18, color=_C_WHITE_P)
        _pptx_txt(slide, str(i + 1), 12.5, 0.12, 0.6, 0.6, size=11, color=_C_WHITE_P, align=PP_ALIGN.RIGHT)
        _pptx_rect(slide, 0.3, 1.0, 12.73, 5.9, _C_WHITE_P)
        if bullets:
            txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.6))
            txb.word_wrap = True
            tf = txb.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(bullets[:12]):
                p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                p.space_before = Pt(4)
                dot = p.add_run()
                dot.text = "\u25cf  "
                dot.font.size = Pt(8)
                dot.font.color.rgb = _C_BLUE_P
                run = p.add_run()
                run.text = bullet
                run.font.size = Pt(12)
                run.font.color.rgb = _C_TXT_P
        else:
            _pptx_txt(slide, "No data available.", 0.5, 1.2, 12, 0.5, size=11, color=_C_DIM_P)
        _footer(slide)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
